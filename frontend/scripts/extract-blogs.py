#!/usr/bin/env python3
"""
extract-blogs.py
Extracts structured country data from PDF/PPTX files in frontend/.blogs
Outputs JSON to frontend/src/data/countries/ and images to frontend/public/images/countries/
"""

import os
import re
import json
import sys
import hashlib
from pathlib import Path

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = Path(__file__).parent
FRONTEND_DIR = SCRIPT_DIR.parent
BLOGS_DIR = FRONTEND_DIR / ".blogs"
JSON_OUT_DIR = FRONTEND_DIR / "src" / "data" / "countries"
IMG_OUT_DIR = FRONTEND_DIR / "public" / "images" / "countries"

JSON_OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── File → Slug Mapping ────────────────────────────────────────────────────
FILE_SLUG_MAP = {
    "Australia - Updated General Presentation.pptx": "australia",
    "Austria.pptx": "austria",
    "Belgium - India.pptx": "belgium",
    "Canada - General - India.pdf": "canada",
    "Cyprus Presentation - India (updated).pdf": "cyprus",
    "Denmark.pptx": "denmark",
    "Dubai General Presentation (Updated).pdf": "dubai",
    "Finland - India.pptx": "finland",
    "France - India.pptx": "france",
    "Georgia - India.pptx": "georgia",
    "Germany - India.pptx": "germany",
    "Hungary - India.pptx": "hungary",
    "Ireland Presentation.pptx": "ireland",
    "Italy.pptx": "italy",
    "JAPAN - India.pptx": "japan",
    "Lithuania - India.pptx": "lithuania",
    "Malaysia Presentation Updated.pdf": "malaysia",
    "Netherlands - India.pptx": "netherlands",
    "New Zealand - India.pptx": "new-zealand",
    "Russia- India.pptx": "russia",
    "Singapore - General presentation - India (Updated).pdf": "singapore",
    "Spain.pptx": "spain",
    "Study in Switzerland - India (Updated).pdf": "switzerland",
    "Sweden - India.pptx": "sweden",
    "UK General Presentation_India-New.pdf": "uk",
    "USA General Presentation.pptx": "usa",
}

# ── Default Unsplash Hero Images per slug ──────────────────────────────────
UNSPLASH_HEROES = {
    "australia": "https://images.unsplash.com/photo-1523482580672-f109ba8cb9be?auto=format&fit=crop&w=1920&q=80",
    "austria": "https://images.unsplash.com/photo-1516550893923-42d28e5677af?auto=format&fit=crop&w=1920&q=80",
    "belgium": "https://images.unsplash.com/photo-1491557345352-5929e343eb89?auto=format&fit=crop&w=1920&q=80",
    "canada": "https://images.unsplash.com/photo-1503614472-8c93d56e92ce?auto=format&fit=crop&w=1920&q=80",
    "cyprus": "https://images.unsplash.com/photo-1541976590-713941681591?auto=format&fit=crop&w=1920&q=80",
    "denmark": "https://images.unsplash.com/photo-1513622470522-26c3c8a854bc?auto=format&fit=crop&w=1920&q=80",
    "dubai": "https://images.unsplash.com/photo-1512453979798-5ea266f8880c?auto=format&fit=crop&w=1920&q=80",
    "finland": "https://images.unsplash.com/photo-1558618666-fcd25c85cd64?auto=format&fit=crop&w=1920&q=80",
    "france": "https://images.unsplash.com/photo-1502602898657-3e91760cbb34?auto=format&fit=crop&w=1920&q=80",
    "georgia": "https://images.unsplash.com/photo-1565008576549-57569a49371d?auto=format&fit=crop&w=1920&q=80",
    "germany": "https://images.unsplash.com/photo-1467269204594-9661b134dd2b?auto=format&fit=crop&w=1920&q=80",
    "hungary": "https://images.unsplash.com/photo-1551867633-194f125bddfa?auto=format&fit=crop&w=1920&q=80",
    "ireland": "https://images.unsplash.com/photo-1590089415225-401ed6f9db8e?auto=format&fit=crop&w=1920&q=80",
    "italy": "https://images.unsplash.com/photo-1515542622106-78bda8ba0e5b?auto=format&fit=crop&w=1920&q=80",
    "japan": "https://images.unsplash.com/photo-1540959733332-eab4deabeeaf?auto=format&fit=crop&w=1920&q=80",
    "lithuania": "https://images.unsplash.com/photo-1559827260-dc66d52bef19?auto=format&fit=crop&w=1920&q=80",
    "malaysia": "https://images.unsplash.com/photo-1596422846543-75c6fc197f07?auto=format&fit=crop&w=1920&q=80",
    "netherlands": "https://images.unsplash.com/photo-1512470876302-972faa2aa9a4?auto=format&fit=crop&w=1920&q=80",
    "new-zealand": "https://images.unsplash.com/photo-1507699622108-4be3abd695ad?auto=format&fit=crop&w=1920&q=80",
    "russia": "https://images.unsplash.com/photo-1513326738677-b964603b136d?auto=format&fit=crop&w=1920&q=80",
    "singapore": "https://images.unsplash.com/photo-1525625293386-3f8f99389edd?auto=format&fit=crop&w=1920&q=80",
    "spain": "https://images.unsplash.com/photo-1543783207-ec64e4d95325?auto=format&fit=crop&w=1920&q=80",
    "switzerland": "https://images.unsplash.com/photo-1527668752968-14dc70a27c95?auto=format&fit=crop&w=1920&q=80",
    "sweden": "https://images.unsplash.com/photo-1509356843151-3e7d96241e11?auto=format&fit=crop&w=1920&q=80",
    "uk": "https://images.unsplash.com/photo-1513635269975-59663e0ac1ad?auto=format&fit=crop&w=1920&q=80",
    "usa": "https://images.unsplash.com/photo-1485738422979-f5c462d49f74?auto=format&fit=crop&w=1920&q=80",
}

COUNTRY_NAMES = {
    "australia": "Australia",
    "austria": "Austria",
    "belgium": "Belgium",
    "canada": "Canada",
    "cyprus": "Cyprus",
    "denmark": "Denmark",
    "dubai": "Dubai (UAE)",
    "finland": "Finland",
    "france": "France",
    "georgia": "Georgia",
    "germany": "Germany",
    "hungary": "Hungary",
    "ireland": "Ireland",
    "italy": "Italy",
    "japan": "Japan",
    "lithuania": "Lithuania",
    "malaysia": "Malaysia",
    "netherlands": "Netherlands",
    "new-zealand": "New Zealand",
    "russia": "Russia",
    "singapore": "Singapore",
    "spain": "Spain",
    "switzerland": "Switzerland",
    "sweden": "Sweden",
    "uk": "United Kingdom",
    "usa": "United States",
}

FLAGS = {
    "australia": "🇦🇺", "austria": "🇦🇹", "belgium": "🇧🇪", "canada": "🇨🇦",
    "cyprus": "🇨🇾", "denmark": "🇩🇰", "dubai": "🇦🇪", "finland": "🇫🇮",
    "france": "🇫🇷", "georgia": "🇬🇪", "germany": "🇩🇪", "hungary": "🇭🇺",
    "ireland": "🇮🇪", "italy": "🇮🇹", "japan": "🇯🇵", "lithuania": "🇱🇹",
    "malaysia": "🇲🇾", "netherlands": "🇳🇱", "new-zealand": "🇳🇿", "russia": "🇷🇺",
    "singapore": "🇸🇬", "spain": "🇪🇸", "switzerland": "🇨🇭", "sweden": "🇸🇪",
    "uk": "🇬🇧", "usa": "🇺🇸",
}

# ── Text Extraction ────────────────────────────────────────────────────────

def extract_text_from_pptx(filepath):
    """Extract all text from a PPTX file, slide by slide."""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            slides_text.append("\n".join(slide_texts))
    return "\n\n--- SLIDE ---\n\n".join(slides_text)


def extract_text_from_pdf(filepath):
    """Extract all text from a PDF file."""
    import fitz
    doc = fitz.open(filepath)
    pages_text = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            pages_text.append(text.strip())
    doc.close()
    return "\n\n--- PAGE ---\n\n".join(pages_text)


def extract_images_from_pptx(filepath, slug):
    """Extract images from a PPTX file and save to public/images/countries/<slug>/"""
    from pptx import Presentation
    from PIL import Image
    import io

    img_dir = IMG_OUT_DIR / slug
    img_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(filepath)
    saved = []
    seen_hashes = set()
    count = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    img_bytes = image.blob
                    img_hash = hashlib.md5(img_bytes).hexdigest()
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)

                    # Filter out tiny images (icons/logos < 50KB)
                    if len(img_bytes) < 30000:
                        continue

                    # Try to open with PIL to get dimensions
                    try:
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        w, h = pil_img.size
                        if w < 200 or h < 100:
                            continue
                        ext = image.ext if image.ext else "jpg"
                        if ext == "jpeg":
                            ext = "jpg"
                        fname = f"img_{count:02d}.{ext}"
                        out_path = img_dir / fname
                        # Save compressed
                        if ext in ("jpg", "jpeg"):
                            pil_img = pil_img.convert("RGB")
                            pil_img.save(out_path, "JPEG", quality=75, optimize=True)
                        else:
                            pil_img.save(out_path)
                        saved.append(f"/images/countries/{slug}/{fname}")
                        count += 1
                    except Exception:
                        pass
                except Exception:
                    pass
    return saved


def extract_images_from_pdf(filepath, slug):
    """Extract images from a PDF file."""
    import fitz
    from PIL import Image
    import io

    img_dir = IMG_OUT_DIR / slug
    img_dir.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(filepath)
    saved = []
    seen_hashes = set()
    count = 0

    for page in doc:
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                base = doc.extract_image(xref)
                img_bytes = base["image"]
                img_hash = hashlib.md5(img_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)

                if len(img_bytes) < 30000:
                    continue

                ext = base.get("ext", "jpg")
                if ext == "jpeg":
                    ext = "jpg"

                try:
                    pil_img = Image.open(io.BytesIO(img_bytes))
                    w, h = pil_img.size
                    if w < 200 or h < 100:
                        continue
                    fname = f"img_{count:02d}.{ext}"
                    out_path = img_dir / fname
                    if ext in ("jpg", "jpeg"):
                        pil_img = pil_img.convert("RGB")
                        pil_img.save(out_path, "JPEG", quality=75, optimize=True)
                    else:
                        pil_img.save(out_path)
                    saved.append(f"/images/countries/{slug}/{fname}")
                    count += 1
                except Exception:
                    pass
            except Exception:
                pass
    doc.close()
    return saved


# ── Text Parsing ───────────────────────────────────────────────────────────

def clean(text):
    """Normalize whitespace."""
    return re.sub(r'\s+', ' ', text).strip()


def find_section(text, keywords, max_chars=800):
    """Find the first paragraph/block containing any of the given keywords."""
    lower = text.lower()
    for kw in keywords:
        idx = lower.find(kw.lower())
        if idx != -1:
            start = max(0, idx - 50)
            snippet = text[start:start + max_chars]
            return clean(snippet)
    return ""


def extract_bullet_lines(text, keywords, num_lines=6):
    """Find lines near a keyword that look like bullet points."""
    lines = text.split("\n")
    results = []
    capture = False
    captured = 0
    for line in lines:
        stripped = line.strip()
        if not stripped:
            if capture and captured > 0:
                capture = False
            continue
        lower = stripped.lower()
        if any(kw.lower() in lower for kw in keywords):
            capture = True
            continue
        if capture:
            if len(stripped) > 10 and captured < num_lines:
                results.append(clean(stripped))
                captured += 1
            elif captured >= num_lines:
                capture = False
    return results


def extract_universities(text):
    """Extract university names from text."""
    uni_keywords = ["university", "college", "institute", "school of", "polytechnic", "hochschule", "fachhochschule", "universität", "université"]
    unis = []
    seen = set()
    for line in text.split("\n"):
        stripped = line.strip()
        if len(stripped) < 5 or len(stripped) > 120:
            continue
        lower = stripped.lower()
        if any(kw in lower for kw in uni_keywords):
            # Clean up common noise
            if any(skip in lower for skip in ["why study", "how to", "benefits", "top ranked", "world class"]):
                continue
            name = clean(stripped).rstrip(".,;:")
            if name and name not in seen and len(name) > 8:
                unis.append(name)
                seen.add(name)
    return unis[:12]


def extract_scholarships(text, slug):
    """Extract scholarship info from text."""
    scholarships = []
    lines = text.split("\n")
    scholarship_kws = ["scholarship", "grant", "bursary", "fellowship", "award", "stipend", "funded", "funding"]
    i = 0
    seen = set()
    while i < len(lines):
        line = lines[i].strip()
        lower = line.lower()
        if any(kw in lower for kw in scholarship_kws) and len(line) > 5 and len(line) < 150:
            name = clean(line).rstrip(".,;:")
            if name not in seen and len(name) > 5:
                seen.add(name)
                # Look ahead for amount/details
                desc_lines = []
                for j in range(i + 1, min(i + 4, len(lines))):
                    next_line = lines[j].strip()
                    if next_line and len(next_line) > 5:
                        desc_lines.append(clean(next_line))
                amount = ""
                eligibility = " ".join(desc_lines[:2]) if desc_lines else "Merit-based scholarship for international students"

                # Try to find monetary value
                money_match = re.search(r'[\$€£¥][\d,]+|[\d,]+\s*(USD|EUR|GBP|AUD|CAD|SGD|CHF|SEK|DKK|HUF|PLN|CZK|JPY)', line + " " + eligibility, re.I)
                if money_match:
                    amount = money_match.group(0)

                scholarships.append({
                    "name": name,
                    "amount": amount or "Varies",
                    "eligibility": eligibility[:200] if eligibility else "Contact institution for details",
                    "link": ""
                })
        i += 1
    return scholarships[:8]


def extract_tuition(text, slug):
    """Extract tuition fee info."""
    currency_map = {
        "australia": "AUD", "austria": "EUR", "belgium": "EUR", "canada": "CAD",
        "cyprus": "EUR", "denmark": "DKK", "dubai": "AED", "finland": "EUR",
        "france": "EUR", "georgia": "USD", "germany": "EUR", "hungary": "EUR",
        "ireland": "EUR", "italy": "EUR", "japan": "JPY", "lithuania": "EUR",
        "malaysia": "MYR", "netherlands": "EUR", "new-zealand": "NZD", "russia": "USD",
        "singapore": "SGD", "spain": "EUR", "switzerland": "CHF", "sweden": "SEK",
        "uk": "GBP", "usa": "USD",
    }
    curr = currency_map.get(slug, "USD")

    # Find tuition-related lines
    lines = text.split("\n")
    tuition_lines = []
    for line in lines:
        stripped = line.strip()
        lower = stripped.lower()
        if any(kw in lower for kw in ["tuition", "fee", "cost", "programme fee", "annual fee", "semester fee"]):
            if re.search(r'[\d,]+', stripped):
                tuition_lines.append(clean(stripped))

    ug = next((l for l in tuition_lines if any(k in l.lower() for k in ["undergrad", "bachelor", "ug", "b.sc", "b.tech", "undergraduate"])), "")
    pg = next((l for l in tuition_lines if any(k in l.lower() for k in ["postgrad", "master", "pg", "mba", "m.sc", "postgraduate", "ms "])), "")
    phd = next((l for l in tuition_lines if any(k in l.lower() for k in ["phd", "doctoral", "research", "doctorate"])), "")

    # Generic fallback from any number found
    if not ug and tuition_lines:
        ug = tuition_lines[0]
    if not pg and len(tuition_lines) > 1:
        pg = tuition_lines[1]

    return {
        "currency": curr,
        "undergraduate": clean(ug) if ug else f"Contact institution for {curr} fees",
        "postgraduate": clean(pg) if pg else f"Contact institution for {curr} fees",
        "phd": clean(phd) if phd else "Often funded through research grants",
    }


def extract_visa(text, slug):
    """Extract visa requirement info."""
    visa_lines = extract_bullet_lines(text, ["visa", "student visa", "visa requirement", "documents required", "required documents"], 8)
    processing = find_section(text, ["processing time", "visa processing", "processing period"], 200)
    fee_match = re.search(r'visa fee[:\s]*([\d,]+\s*(?:USD|EUR|GBP|AUD|INR|€|\$|£)?)', text, re.I)
    visa_type_map = {
        "australia": "Student Visa (Subclass 500)",
        "austria": "Student Visa (Type D)",
        "belgium": "Student Visa Type D",
        "canada": "Canadian Student Permit",
        "cyprus": "Student Visa Category D",
        "denmark": "Danish Residence Permit (Student)",
        "dubai": "UAE Student Visa",
        "finland": "Finnish Residence Permit (Student)",
        "france": "Long-Stay Student Visa (VLS-TS)",
        "georgia": "Georgian Student Visa",
        "germany": "German Student Visa (Section 16b AufenthG)",
        "hungary": "Hungarian Student Visa (D Visa)",
        "ireland": "Irish Study Visa",
        "italy": "Italian Student Visa (Type D)",
        "japan": "Japanese Student Visa (College Student)",
        "lithuania": "Lithuanian National Visa (D)",
        "malaysia": "Malaysia Student Pass",
        "netherlands": "Dutch MVV / Residence Permit (Student)",
        "new-zealand": "New Zealand Student Visa",
        "russia": "Russian Student Visa",
        "singapore": "Singapore Student Pass",
        "spain": "Spanish Student Visa (Type D)",
        "switzerland": "Swiss Student Visa (Type D)",
        "sweden": "Swedish Residence Permit (Student)",
        "uk": "UK Student Visa (formerly Tier 4)",
        "usa": "F-1 Student Visa",
    }

    reqs = visa_lines if visa_lines else [
        "Valid passport (6+ months validity)",
        "Unconditional offer letter from university",
        "Proof of financial sufficiency",
        "English language test scores",
        "Medical insurance",
        "Passport-sized photographs",
    ]

    return {
        "type": visa_type_map.get(slug, "Student Visa"),
        "requirements": reqs,
        "processingTime": clean(processing[:150]) if processing else "4–8 weeks",
        "fee": fee_match.group(1).strip() if fee_match else "Varies by country of application",
    }


def extract_work_permit(text, slug):
    """Extract work permit / post-study work rights."""
    during_lines = find_section(text, ["part-time", "work during", "hours per week", "20 hours", "study permit work"], 300)
    post_lines = find_section(text, ["post study", "post-study", "stay back", "graduate visa", "work permit after", "pgwp", "opt", "psw", "graduate route"], 300)

    defaults = {
        "australia": {"duringStudy": "Up to 48 hours/fortnight during semester, unlimited during breaks", "postStudy": "2–4 years (Graduate Visa subclass 485)", "notes": "Duration depends on level and location of study"},
        "austria": {"duringStudy": "Up to 20 hours/week", "postStudy": "12 months job-seeker visa", "notes": "Red-White-Red card available post-graduation"},
        "belgium": {"duringStudy": "Up to 20 hours/week during term, full-time in holidays", "postStudy": "12-month job-search permit for graduates", "notes": ""},
        "canada": {"duringStudy": "Up to 24 hours/week off-campus; unlimited on-campus", "postStudy": "Up to 3 years PGWP (equal to study duration)", "notes": "PGWP is a direct pathway to Canadian PR"},
        "cyprus": {"duringStudy": "Up to 20 hours/week", "postStudy": "Temporary residence extension available", "notes": ""},
        "denmark": {"duringStudy": "Up to 20 hours/week; full-time in June, July, August", "postStudy": "2-year job-seeker permit", "notes": ""},
        "dubai": {"duringStudy": "Part-time work permitted with NOC from institution", "postStudy": "2-year graduate visa available", "notes": ""},
        "finland": {"duringStudy": "Up to 25 hours/week during term", "postStudy": "Up to 2 years extended permit for job-seeking", "notes": ""},
        "france": {"duringStudy": "Up to 964 hours/year (approx. 20 hours/week)", "postStudy": "APS Visa: 1 year to find a job", "notes": ""},
        "georgia": {"duringStudy": "No official restriction for student workers", "postStudy": "Temporary residence permit available", "notes": "Low-cost destination with growing tech sector"},
        "germany": {"duringStudy": "120 full days or 240 half days per year", "postStudy": "18-month job-seeker visa", "notes": "Opportunity Card (Chancenkarte) for skilled workers"},
        "hungary": {"duringStudy": "Up to 24 hours/week during semester", "postStudy": "9-month job-seeker permit", "notes": ""},
        "ireland": {"duringStudy": "Up to 20 hours/week during term; 40 hours during summer", "postStudy": "1–2 years Stay Back Option (SBOT)", "notes": "Degree graduates get 2 years; sub-degree get 1 year"},
        "italy": {"duringStudy": "Up to 20 hours/week", "postStudy": "12-month job-seeker stay", "notes": ""},
        "japan": {"duringStudy": "Up to 28 hours/week with permission", "postStudy": "Designated Activities visa for job-hunting (up to 1 year)", "notes": ""},
        "lithuania": {"duringStudy": "Up to 20 hours/week during term", "postStudy": "Temporary residence permit for 6–12 months", "notes": ""},
        "malaysia": {"duringStudy": "Up to 20 hours/week during holidays", "postStudy": "Graduate Employment Pass available", "notes": ""},
        "netherlands": {"duringStudy": "Up to 16 hours/week; full-time in June, July, August", "postStudy": "1-year Orientation Year visa (Zoekjaar)", "notes": ""},
        "new-zealand": {"duringStudy": "Up to 20 hours/week; full-time in scheduled breaks", "postStudy": "1–3 years Post-Study Work Visa", "notes": ""},
        "russia": {"duringStudy": "Part-time work with university permission", "postStudy": "Temporary work permit available", "notes": "Low tuition and cost of living"},
        "singapore": {"duringStudy": "Up to 16 hours/week during term", "postStudy": "Employment Pass / S Pass for graduates", "notes": ""},
        "spain": {"duringStudy": "Up to 30 hours/week", "postStudy": "1-year job-seeker extension available", "notes": ""},
        "switzerland": {"duringStudy": "Up to 15 hours/week during term", "postStudy": "6-month job-seeker permit post-graduation", "notes": ""},
        "sweden": {"duringStudy": "Unlimited hours during study period", "postStudy": "1-year job-seeker permit", "notes": ""},
        "uk": {"duringStudy": "Up to 20 hours/week during term; full-time in official breaks", "postStudy": "2–3 years Graduate Route Visa (PSW)", "notes": "STEM graduates may qualify for 3 years"},
        "usa": {"duringStudy": "Up to 20 hours/week on-campus; OPT during study", "postStudy": "1 year OPT; 3 years for STEM graduates", "notes": "H-1B lottery for long-term US employment"},
    }

    d = defaults.get(slug, {})
    return {
        "duringStudy": clean(during_lines[:200]) if during_lines else d.get("duringStudy", "Part-time work permitted"),
        "postStudy": clean(post_lines[:200]) if post_lines else d.get("postStudy", "Post-study work permit available"),
        "notes": d.get("notes", ""),
    }


def extract_language_requirements(text, slug):
    """Extract IELTS/TOEFL/PTE requirements."""
    ielts_match = re.search(r'IELTS[:\s]*([0-9.]+)', text, re.I)
    toefl_match = re.search(r'TOEFL[:\s]*(?:iBT)?[:\s]*([0-9]+)', text, re.I)
    pte_match = re.search(r'PTE[:\s]*(?:Academic)?[:\s]*([0-9]+)', text, re.I)
    gmat_match = re.search(r'GMAT[:\s]*([0-9]+)', text, re.I)
    gre_match = re.search(r'GRE[:\s]*([0-9]+)', text, re.I)

    # Defaults for non-English speaking countries
    non_english = ["austria", "belgium", "denmark", "finland", "france", "germany", "hungary", "italy", "japan", "lithuania", "netherlands", "russia", "spain", "switzerland", "sweden", "georgia", "cyprus", "malaysia", "singapore"]

    notes_map = {
        "germany": "Some programs taught in English; German B2 required for German-taught programs",
        "france": "French B2 required for French-taught programs; many English-taught programs available",
        "austria": "German B2 for German programs; English programs widely available",
        "netherlands": "Many programs taught in English; Dutch not required",
        "sweden": "Most master's programs taught in English",
        "finland": "Many programs taught in English; Finnish required for Finnish-taught programs",
        "denmark": "Many programs taught in English; Danish required for some programs",
        "switzerland": "French, German, or Italian depending on region; many English programs at top schools",
        "belgium": "French or Dutch depending on region; many English programs available",
        "japan": "Japanese JLPT N2 for Japanese-taught programs; English programs growing rapidly",
        "russia": "Russian language proficiency for Russian-taught programs; English programs available",
        "spain": "Spanish B1–B2 for Spanish programs; English programs at private universities",
        "italy": "Italian B2 for Italian programs; growing number of English programs",
        "hungary": "Some programs in English; Hungarian for local programs",
        "georgia": "English proficiency required for English programs; Georgian not required",
        "lithuania": "English programs widely available; Lithuanian for local programs",
        "malaysia": "English medium instruction; IELTS 6.0+ for top universities",
        "singapore": "English is the primary medium of instruction",
        "cyprus": "Many programs taught in English or Greek",
        "dubai": "English medium instruction widely available",
    }

    return {
        "ielts": ielts_match.group(1) if ielts_match else ("6.0" if slug in non_english else "6.5"),
        "toefl": toefl_match.group(1) if toefl_match else ("80" if slug in non_english else "90"),
        "pte": pte_match.group(1) if pte_match else ("58" if slug in non_english else "65"),
        "gmat": gmat_match.group(1) if gmat_match else "",
        "gre": gre_match.group(1) if gre_match else "",
        "notes": notes_map.get(slug, "English language proficiency required for all programs"),
    }


def extract_popular_programs(text, slug):
    """Extract popular study programs."""
    program_keywords = [
        "business administration", "mba", "computer science", "information technology",
        "engineering", "medicine", "law", "economics", "finance", "architecture",
        "design", "arts", "psychology", "public health", "data science", "ai",
        "artificial intelligence", "machine learning", "biotechnology", "pharmacy",
        "nursing", "hotel management", "hospitality", "fashion", "media", "marketing",
        "accounting", "management", "environmental", "renewable energy", "cybersecurity",
        "software engineering", "electrical engineering", "mechanical engineering",
        "civil engineering", "chemical engineering", "mathematics", "physics",
        "biology", "chemistry", "international relations", "political science",
    ]

    found = []
    lower_text = text.lower()
    for prog in program_keywords:
        if prog in lower_text:
            found.append(prog.title())

    # Country-specific defaults
    defaults = {
        "australia": ["Engineering", "Business Administration", "Information Technology", "Medicine", "Nursing", "Environmental Science"],
        "austria": ["Music & Arts", "Architecture", "Business Administration", "Engineering", "Tourism Management", "Medicine"],
        "belgium": ["Business Administration", "European Studies", "Law", "Economics", "Engineering", "Arts"],
        "canada": ["Computer Science", "Engineering", "Business Administration", "Medicine", "Data Science", "Environmental Science"],
        "cyprus": ["Business Administration", "Law", "Computer Science", "Hospitality Management", "Architecture", "Engineering"],
        "denmark": ["Design", "Architecture", "Engineering", "Business Administration", "Life Sciences", "Environmental Studies"],
        "dubai": ["Business Administration", "Engineering", "Hospitality Management", "Finance", "Law", "Computer Science"],
        "finland": ["Engineering", "Computer Science", "Design", "Business Administration", "Environmental Science", "Education"],
        "france": ["Fashion & Design", "Business Administration", "Engineering", "Culinary Arts", "International Relations", "Law"],
        "georgia": ["Medicine", "Dentistry", "Business Administration", "Engineering", "Law", "Computer Science"],
        "germany": ["Engineering", "Computer Science", "Business Administration", "Automotive Engineering", "Medicine", "Physics"],
        "hungary": ["Medicine", "Dentistry", "Engineering", "Business Administration", "Computer Science", "Pharmacy"],
        "ireland": ["Computer Science", "Pharmacy", "Business Administration", "Data Analytics", "Engineering", "Finance"],
        "italy": ["Architecture", "Fashion Design", "Arts", "Business Administration", "Law", "Engineering"],
        "japan": ["Engineering", "Robotics", "Business Administration", "Arts", "Computer Science", "Medicine"],
        "lithuania": ["Medicine", "Dentistry", "Engineering", "Business Administration", "Computer Science", "Law"],
        "malaysia": ["Engineering", "Business Administration", "Computer Science", "Medicine", "Hospitality", "Finance"],
        "netherlands": ["Business Administration", "Engineering", "Law", "Economics", "International Relations", "Computer Science"],
        "new-zealand": ["Agriculture", "Engineering", "Business Administration", "Environmental Science", "Computer Science", "Tourism"],
        "russia": ["Medicine", "Engineering", "Computer Science", "Physics", "Mathematics", "Arts"],
        "singapore": ["Finance", "Business Administration", "Engineering", "Computer Science", "Medicine", "Law"],
        "spain": ["Business Administration", "Tourism", "Architecture", "Law", "Engineering", "Arts"],
        "switzerland": ["Finance", "International Relations", "Business Administration", "Engineering", "Hospitality", "Medicine"],
        "sweden": ["Engineering", "Computer Science", "Design", "Business Administration", "Environmental Science", "Medicine"],
        "uk": ["Business Administration", "Law", "Engineering", "Computer Science", "Medicine", "Finance"],
        "usa": ["Computer Science", "Business Administration", "Engineering", "Medicine", "Data Science", "Finance"],
    }

    if not found:
        return defaults.get(slug, ["Business Administration", "Engineering", "Computer Science", "Medicine", "Law"])
    return list(dict.fromkeys(found + defaults.get(slug, [])))[:8]


def extract_faqs(text, slug):
    """Generate or extract FAQs."""
    faq_defaults = {
        "australia": [
            {"question": "Can Indian students work while studying in Australia?", "answer": "Yes, student visa holders can work up to 48 hours per fortnight during the academic semester and unlimited hours during scheduled breaks."},
            {"question": "What is the PSW visa in Australia?", "answer": "The Temporary Graduate Visa (subclass 485) allows international graduates to live and work in Australia for 2–4 years after completing their studies."},
            {"question": "Is IELTS mandatory for Australian universities?", "answer": "Most universities require an IELTS score of 6.0–6.5 (overall). Some accept alternative tests like PTE, TOEFL, or Duolingo."},
            {"question": "What are the popular intakes in Australia?", "answer": "The main intakes are February (Semester 1) and July (Semester 2). Some institutions also offer November intake."},
        ],
        "canada": [
            {"question": "What is the PGWP in Canada?", "answer": "The Post-Graduation Work Permit (PGWP) allows graduates to work in Canada for up to 3 years. It is a major pathway to Canadian Permanent Residency."},
            {"question": "Can I bring my family to Canada as a student?", "answer": "Yes, spouses of international students can apply for an open work permit, and children can attend school in Canada."},
            {"question": "How long does it take to get a Canadian student permit?", "answer": "Processing typically takes 4–8 weeks for online applications. Biometrics may be required for certain nationalities."},
            {"question": "What is the SDS (Student Direct Stream) in Canada?", "answer": "SDS is a faster processing stream for students from eligible countries including India. Applications are usually processed within 20 working days."},
        ],
        "uk": [
            {"question": "What is the UK Graduate Route Visa?", "answer": "The Graduate Route allows international students to stay in the UK for 2 years (3 for PhD graduates) after completing their degree to work or look for work."},
            {"question": "Do I need a CAS number for a UK Student Visa?", "answer": "Yes, you need a Confirmation of Acceptance for Studies (CAS) from your UK university before applying for a Student Visa."},
            {"question": "How much bank balance is needed for a UK Student Visa?", "answer": "You must show £1,334/month for up to 9 months for London, or £1,023/month for outside London, plus your full first-year tuition fees."},
            {"question": "Can I work while studying in the UK?", "answer": "Yes, most student visa holders can work up to 20 hours per week during term time and full-time during official vacation periods."},
        ],
        "germany": [
            {"question": "Is studying in Germany really free?", "answer": "Public universities in Germany charge only a semester fee of €150–€350 for administration. There is no tuition fee for most programs, even for international students."},
            {"question": "Do I need to know German to study in Germany?", "answer": "Not necessarily. Many master's programs are offered in English. However, knowing German is very helpful for daily life and increases job prospects."},
            {"question": "What is the blocked account requirement for Germany?", "answer": "A blocked account (Sperrkonto) of €11,904/year (€992/month) is required to show financial sufficiency for your German student visa."},
            {"question": "Can I stay in Germany after graduation?", "answer": "Yes, you can apply for an 18-month job-seeker visa (Aufenthaltserlaubnis zur Arbeitssuche) after completing your degree."},
        ],
        "usa": [
            {"question": "What is F-1 OPT?", "answer": "Optional Practical Training (OPT) allows F-1 students to work in a job related to their major for up to 12 months. STEM graduates can apply for a 24-month extension."},
            {"question": "How much does it cost to study in the USA?", "answer": "Tuition ranges from $20,000–$60,000/year at private universities and $10,000–$30,000/year at public universities for out-of-state students."},
            {"question": "What are the SAT/GRE/GMAT requirements?", "answer": "Requirements vary by program. Most graduate programs require GRE or GMAT. Undergraduate admissions often consider SAT/ACT scores."},
            {"question": "Can Indian students get scholarships in the USA?", "answer": "Yes. Many universities offer merit scholarships, and there are federal programs like Fulbright. Assistantships (TA/RA) also provide stipends for graduate students."},
        ],
    }

    return faq_defaults.get(slug, [
        {"question": f"Is it easy to get a student visa for {COUNTRY_NAMES.get(slug, slug)}?", "answer": f"With proper documentation and financial proof, obtaining a student visa for {COUNTRY_NAMES.get(slug, slug)} is straightforward. MVR Consultants can guide you through the entire process."},
        {"question": f"What are the tuition fees in {COUNTRY_NAMES.get(slug, slug)}?", "answer": "Tuition fees vary by university and program. Contact our counselors for program-specific fee structures."},
        {"question": "Can I work while studying?", "answer": "Most countries allow international students to work part-time during term time. Specific hours and conditions depend on the country's immigration rules."},
        {"question": "Are there scholarships available for Indian students?", "answer": "Yes, many universities and governments offer merit-based and need-based scholarships for international students. Ask our counselors for scholarship matching."},
    ])


def build_country_json(slug, text, images):
    """Build the structured JSON for a country."""
    name = COUNTRY_NAMES.get(slug, slug.title())

    unis = extract_universities(text)
    scholarships = extract_scholarships(text, slug)
    tuition = extract_tuition(text, slug)
    visa = extract_visa(text, slug)
    work_permit = extract_work_permit(text, slug)
    lang_req = extract_language_requirements(text, slug)
    programs = extract_popular_programs(text, slug)
    faqs = extract_faqs(text, slug)

    # Stats
    stats_defaults = {
        "australia": {"unis": "40+", "workRights": "Up to 4 Years", "intakes": "Feb, Jul", "avgTuition": "AUD 30k–50k/yr"},
        "austria": {"unis": "70+", "workRights": "12 Months", "intakes": "Oct, Mar", "avgTuition": "€700–1,500/yr"},
        "belgium": {"unis": "50+", "workRights": "12 Months", "intakes": "Sep, Feb", "avgTuition": "€835–4,175/yr"},
        "canada": {"unis": "100+", "workRights": "Up to 3 Years", "intakes": "Sep, Jan, May", "avgTuition": "CAD 15k–35k/yr"},
        "cyprus": {"unis": "10+", "workRights": "Temp. Permit", "intakes": "Sep, Feb", "avgTuition": "€3,500–8,000/yr"},
        "denmark": {"unis": "40+", "workRights": "2 Years", "intakes": "Sep, Feb", "avgTuition": "DKK 45k–120k/yr"},
        "dubai": {"unis": "60+", "workRights": "2 Years", "intakes": "Sep, Jan, May", "avgTuition": "AED 30k–80k/yr"},
        "finland": {"unis": "40+", "workRights": "2 Years", "intakes": "Sep", "avgTuition": "€6,000–18,000/yr"},
        "france": {"unis": "250+", "workRights": "1 Year (APS)", "intakes": "Sep, Jan", "avgTuition": "€2,770–16,000/yr"},
        "georgia": {"unis": "50+", "workRights": "Temp. Permit", "intakes": "Sep, Feb", "avgTuition": "$2,000–6,000/yr"},
        "germany": {"unis": "300+", "workRights": "18 Months", "intakes": "Oct, Apr", "avgTuition": "€0–3,000/yr"},
        "hungary": {"unis": "65+", "workRights": "9 Months", "intakes": "Sep, Feb", "avgTuition": "€3,000–14,000/yr"},
        "ireland": {"unis": "30+", "workRights": "Up to 2 Years", "intakes": "Sep, Jan", "avgTuition": "€10k–25k/yr"},
        "italy": {"unis": "90+", "workRights": "12 Months", "intakes": "Oct, Mar", "avgTuition": "€900–12,000/yr"},
        "japan": {"unis": "750+", "workRights": "1 Year", "intakes": "Apr, Sep", "avgTuition": "¥500k–2M/yr"},
        "lithuania": {"unis": "20+", "workRights": "6–12 Months", "intakes": "Sep, Feb", "avgTuition": "€2,500–5,000/yr"},
        "malaysia": {"unis": "70+", "workRights": "Emp. Pass", "intakes": "Mar, Jul, Nov", "avgTuition": "MYR 15k–40k/yr"},
        "netherlands": {"unis": "50+", "workRights": "1 Year", "intakes": "Sep, Feb", "avgTuition": "€8,000–20,000/yr"},
        "new-zealand": {"unis": "8+", "workRights": "Up to 3 Years", "intakes": "Feb, Jul", "avgTuition": "NZD 22k–35k/yr"},
        "russia": {"unis": "300+", "workRights": "Temp. Permit", "intakes": "Sep, Feb", "avgTuition": "$2,000–6,000/yr"},
        "singapore": {"unis": "6+", "workRights": "Emp. Pass", "intakes": "Aug, Jan", "avgTuition": "SGD 15k–35k/yr"},
        "spain": {"unis": "80+", "workRights": "1 Year", "intakes": "Sep, Feb", "avgTuition": "€1,000–12,000/yr"},
        "switzerland": {"unis": "20+", "workRights": "6 Months", "intakes": "Sep, Feb", "avgTuition": "CHF 1,000–2,000/yr"},
        "sweden": {"unis": "40+", "workRights": "1 Year", "intakes": "Sep, Jan", "avgTuition": "SEK 80k–180k/yr"},
        "uk": {"unis": "150+", "workRights": "2–3 Years", "intakes": "Sep, Jan", "avgTuition": "£12k–25k/yr"},
        "usa": {"unis": "4,000+", "workRights": "1–3 Years (OPT)", "intakes": "Fall, Spring", "avgTuition": "$20k–60k/yr"},
    }

    stats = stats_defaults.get(slug, {"unis": "50+", "workRights": "Varies", "intakes": "Sep, Jan", "avgTuition": "Contact for fees"})

    taglines = {
        "australia": "Global Opportunities", "austria": "Alpine Academic Excellence",
        "belgium": "Heart of Europe", "canada": "Welcoming & Diverse",
        "cyprus": "Mediterranean Study Hub", "denmark": "Innovation & Design",
        "dubai": "Future-Forward Education", "finland": "Nordic Innovation",
        "france": "Culture Meets Excellence", "georgia": "Affordable Gateway to Europe",
        "germany": "Engineering & Tech Hub", "hungary": "Europe's Medical Capital",
        "ireland": "Europe's Silicon Valley", "italy": "Art, Culture & Academia",
        "japan": "Innovation Meets Tradition", "lithuania": "Baltic Academic Hub",
        "malaysia": "Asia's Education Hub", "netherlands": "Bicycle & Business Capital",
        "new-zealand": "Nature & Knowledge", "russia": "Science & Heritage",
        "singapore": "Asia's Global City", "spain": "Sun, Culture & Study",
        "switzerland": "Precision & Excellence", "sweden": "Nordic Tech Pioneer",
        "uk": "Academic Excellence", "usa": "Hub of Innovation",
    }

    descriptions = {
        "australia": "Home to 8 of the world's top 100 universities, Australia combines world-class education with an incredible lifestyle. Strong post-study work options and a clear pathway to permanent residency make it a top choice for Indian students.",
        "austria": "Austria offers a rich academic tradition, affordable tuition, and a central European location. Vienna is consistently ranked among the world's most livable cities, making it perfect for students.",
        "belgium": "Located at the heart of Europe, Belgium is home to some of the continent's oldest universities and serves as the political capital of the EU. Excellent for students interested in international relations, law, and business.",
        "canada": "Consistently ranked as one of the best countries for quality of life, Canada offers top-tier education with clear pathways to permanent residency. The PGWP makes it a top choice for Indian students seeking long-term immigration.",
        "cyprus": "A sunny Mediterranean island nation, Cyprus offers affordable, quality education in an English-friendly environment. Its EU membership and strategic location make it an attractive base for students.",
        "denmark": "Denmark is a global leader in design, sustainability, and innovation. Its universities consistently rank among Europe's best, and the country boasts a high quality of life with a strong focus on work-life balance.",
        "dubai": "Dubai and the UAE offer a tax-free, cosmopolitan education hub with campuses of top global universities. The region's booming economy creates excellent career opportunities for graduates.",
        "finland": "Finland's education system is among the world's best. Top-ranked universities, a strong focus on research and innovation, and a growing number of English-taught programs make Finland increasingly popular.",
        "france": "France is a global hub of culture, fashion, and innovation. Home to world-class Grandes Écoles and public universities, France offers excellent education at remarkably affordable tuition rates.",
        "georgia": "Georgia is an emerging affordable destination for medical and engineering students. Its universities are WHO-recognized and MCI-approved, making it an excellent choice for Indian medical aspirants.",
        "germany": "The engine of Europe's economy. Germany is globally revered for engineering and technology, offering zero or extremely low tuition fees at public universities — even for international students.",
        "hungary": "Hungary is a top destination for medical education in Europe, with WHO-recognized programs and affordable fees. Its vibrant capital Budapest offers a rich student life.",
        "ireland": "The European headquarters for Google, Apple, and Facebook. Ireland offers an English-speaking environment with a booming tech and pharma sector, making it ideal for career-focused students.",
        "italy": "Italy blends world-class education with unparalleled cultural heritage. Home to some of the world's oldest universities, it offers programs in architecture, arts, fashion, and engineering.",
        "japan": "A global leader in technology and innovation, Japan offers cutting-edge research opportunities and a unique cultural experience. Government scholarships like MEXT make it accessible for international students.",
        "lithuania": "Lithuania is an affordable EU destination with WHO-recognized medical programs and a growing tech sector. English-taught programs are widely available at accredited universities.",
        "malaysia": "Malaysia is Asia's most accessible study destination, offering quality education at affordable costs. Its multicultural environment and English medium instruction make international students feel at home.",
        "netherlands": "The Netherlands is a forward-thinking nation with world-class universities. It offers hundreds of English-taught bachelor's and master's programs, making it a top choice for international students in Europe.",
        "new-zealand": "New Zealand offers a safe, welcoming environment with world-class universities and a stunning natural landscape. Its practical approach to education and generous work rights attract students globally.",
        "russia": "Russia is home to world-leading universities in science, engineering, and medicine. Scholarships through the Russian Government Scholarship program make it a highly accessible option.",
        "singapore": "Singapore is Asia's global financial capital with world-class universities like NUS and NTU. Its strategic location, safe environment, and strong economy make it a premium study destination.",
        "spain": "Spain offers an excellent quality of life, rich culture, and affordable living costs. Its universities provide strong programs in business, architecture, and tourism within an EU framework.",
        "switzerland": "Switzerland is home to some of the world's top universities including ETH Zurich. Despite low tuition fees at public universities, the high cost of living in Switzerland requires careful financial planning.",
        "sweden": "Sweden is a global innovator with top-ranked universities and a strong focus on sustainability and technology. Many programs are English-taught, and the country offers a high standard of living.",
        "uk": "Experience academic prestige with shorter course durations. The UK is a global leader in quality education with the 2–3 year Graduate Route Visa making it one of the most sought-after study destinations.",
        "usa": "The US is home to the world's highest number of top-ranked universities, offering unparalleled academic flexibility, research opportunities, and career prospects through OPT and STEM extensions.",
    }

    top_unis_map = {
        "australia": [{"name": "University of Melbourne", "ranking": "#33 QS 2024", "programs": ["Engineering", "Business", "Medicine"]}, {"name": "Australian National University", "ranking": "#34 QS 2024", "programs": ["Sciences", "Arts", "Law"]}, {"name": "University of Sydney", "ranking": "#42 QS 2024", "programs": ["Medicine", "Engineering", "Law"]}, {"name": "UNSW Sydney", "ranking": "#45 QS 2024", "programs": ["Engineering", "Business", "Law"]}],
        "austria": [{"name": "University of Vienna", "ranking": "#150 QS 2024", "programs": ["Law", "Sciences", "Humanities"]}, {"name": "TU Wien", "ranking": "#Top 300 QS", "programs": ["Engineering", "Technology", "Architecture"]}, {"name": "Vienna University of Economics", "ranking": "AACSB Accredited", "programs": ["Business", "Economics", "Finance"]}],
        "belgium": [{"name": "KU Leuven", "ranking": "#74 QS 2024", "programs": ["Engineering", "Medicine", "Sciences"]}, {"name": "Ghent University", "ranking": "#Top 150 QS", "programs": ["Sciences", "Engineering", "Veterinary"]}, {"name": "Université Libre de Bruxelles", "ranking": "Top 200 QS", "programs": ["Law", "Sciences", "Business"]}],
        "canada": [{"name": "University of Toronto", "ranking": "#21 QS 2024", "programs": ["Engineering", "Medicine", "Business"]}, {"name": "McGill University", "ranking": "#30 QS 2024", "programs": ["Medicine", "Law", "Sciences"]}, {"name": "University of British Columbia", "ranking": "#34 QS 2024", "programs": ["Engineering", "Business", "Sciences"]}, {"name": "University of Waterloo", "ranking": "#112 QS 2024", "programs": ["Computer Science", "Engineering", "Math"]}],
        "cyprus": [{"name": "University of Cyprus", "ranking": "Top 600 QS", "programs": ["Business", "Engineering", "Humanities"]}, {"name": "Cyprus University of Technology", "ranking": "Nationally Accredited", "programs": ["Engineering", "IT", "Business"]}, {"name": "European University Cyprus", "ranking": "Accredited", "programs": ["Medicine", "Dentistry", "Business"]}],
        "denmark": [{"name": "Technical University of Denmark (DTU)", "ranking": "Top 200 QS", "programs": ["Engineering", "Life Sciences", "IT"]}, {"name": "University of Copenhagen", "ranking": "#79 QS 2024", "programs": ["Sciences", "Medicine", "Humanities"]}, {"name": "Copenhagen Business School", "ranking": "AACSB Accredited", "programs": ["Business", "Finance", "International Relations"]}],
        "dubai": [{"name": "University of Dubai", "ranking": "Nationally Accredited", "programs": ["Business", "Law", "Engineering"]}, {"name": "Middlesex University Dubai", "ranking": "UK Accredited", "programs": ["Business", "Design", "Engineering"]}, {"name": "Heriot-Watt University Dubai", "ranking": "UK Accredited", "programs": ["Engineering", "Business", "Computer Science"]}],
        "finland": [{"name": "University of Helsinki", "ranking": "#106 QS 2024", "programs": ["Sciences", "Law", "Humanities"]}, {"name": "Aalto University", "ranking": "#Top 150 QS", "programs": ["Engineering", "Business", "Design"]}, {"name": "Tampere University", "ranking": "Top 400 QS", "programs": ["Engineering", "IT", "Health Sciences"]}],
        "france": [{"name": "Sorbonne University", "ranking": "#59 QS 2024", "programs": ["Sciences", "Humanities", "Medicine"]}, {"name": "École Polytechnique", "ranking": "#61 QS 2024", "programs": ["Engineering", "Sciences", "Economics"]}, {"name": "HEC Paris", "ranking": "#1 European Business School", "programs": ["MBA", "Business", "Finance"]}, {"name": "Sciences Po Paris", "ranking": "Top 250 QS", "programs": ["Political Science", "Law", "International Relations"]}],
        "georgia": [{"name": "Tbilisi State Medical University", "ranking": "WHO Recognized", "programs": ["Medicine", "Dentistry", "Pharmacy"]}, {"name": "University of Georgia (Tbilisi)", "ranking": "MCI Approved", "programs": ["Medicine", "Business", "Law"]}, {"name": "Caucasus University", "ranking": "Nationally Accredited", "programs": ["Business", "Law", "IT"]}],
        "germany": [{"name": "Technical University of Munich", "ranking": "#37 QS 2024", "programs": ["Engineering", "Computer Science", "Sciences"]}, {"name": "LMU Munich", "ranking": "#54 QS 2024", "programs": ["Medicine", "Law", "Sciences"]}, {"name": "Heidelberg University", "ranking": "#87 QS 2024", "programs": ["Medicine", "Sciences", "Humanities"]}, {"name": "RWTH Aachen", "ranking": "Top 100 QS", "programs": ["Engineering", "Technology", "Business"]}],
        "hungary": [{"name": "Semmelweis University", "ranking": "Top Medical School", "programs": ["Medicine", "Dentistry", "Pharmacy"]}, {"name": "Budapest University of Technology", "ranking": "Top 300 QS", "programs": ["Engineering", "IT", "Architecture"]}, {"name": "Corvinus University", "ranking": "EQUIS Accredited", "programs": ["Business", "Economics", "Law"]}],
        "ireland": [{"name": "Trinity College Dublin", "ranking": "#81 QS 2024", "programs": ["Computer Science", "Law", "Medicine"]}, {"name": "University College Dublin", "ranking": "#181 QS 2024", "programs": ["Business", "Engineering", "Sciences"]}, {"name": "University of Galway", "ranking": "Top 300 QS", "programs": ["Medicine", "Arts", "Sciences"]}, {"name": "University College Cork", "ranking": "Top 300 QS", "programs": ["Medicine", "Business", "Engineering"]}],
        "italy": [{"name": "University of Bologna", "ranking": "#154 QS 2024", "programs": ["Law", "Engineering", "Arts"]}, {"name": "Politecnico di Milano", "ranking": "#139 QS 2024", "programs": ["Architecture", "Engineering", "Design"]}, {"name": "Sapienza University of Rome", "ranking": "Top 200 QS", "programs": ["Medicine", "Law", "Sciences"]}],
        "japan": [{"name": "University of Tokyo", "ranking": "#28 QS 2024", "programs": ["Engineering", "Sciences", "Law"]}, {"name": "Kyoto University", "ranking": "#46 QS 2024", "programs": ["Sciences", "Medicine", "Engineering"]}, {"name": "Osaka University", "ranking": "#80 QS 2024", "programs": ["Engineering", "Medicine", "Sciences"]}],
        "lithuania": [{"name": "Vilnius University", "ranking": "Top 600 QS", "programs": ["Medicine", "Law", "IT"]}, {"name": "Lithuanian University of Health Sciences", "ranking": "WHO Recognized", "programs": ["Medicine", "Dentistry", "Pharmacy"]}, {"name": "KTU (Kaunas University of Technology)", "ranking": "Nationally Accredited", "programs": ["Engineering", "IT", "Business"]}],
        "malaysia": [{"name": "University of Malaya", "ranking": "#65 QS 2024", "programs": ["Medicine", "Engineering", "Business"]}, {"name": "Universiti Putra Malaysia", "ranking": "Top 200 QS", "programs": ["Agriculture", "Sciences", "Engineering"]}, {"name": "Taylor's University", "ranking": "Top 250 QS", "programs": ["Hospitality", "Business", "Engineering"]}],
        "netherlands": [{"name": "Delft University of Technology", "ranking": "#47 QS 2024", "programs": ["Engineering", "Architecture", "Computer Science"]}, {"name": "University of Amsterdam", "ranking": "#53 QS 2024", "programs": ["Business", "Sciences", "Humanities"]}, {"name": "Eindhoven University of Technology", "ranking": "Top 150 QS", "programs": ["Engineering", "Technology", "Business"]}],
        "new-zealand": [{"name": "University of Auckland", "ranking": "#65 QS 2024", "programs": ["Engineering", "Medicine", "Business"]}, {"name": "University of Otago", "ranking": "Top 250 QS", "programs": ["Medicine", "Sciences", "Business"]}, {"name": "Victoria University of Wellington", "ranking": "Top 250 QS", "programs": ["Law", "Business", "Sciences"]}],
        "russia": [{"name": "Lomonosov Moscow State University", "ranking": "#87 QS 2024", "programs": ["Sciences", "Mathematics", "Engineering"]}, {"name": "Saint Petersburg State University", "ranking": "Top 250 QS", "programs": ["Sciences", "Law", "Humanities"]}, {"name": "ITMO University", "ranking": "Top IT University", "programs": ["Computer Science", "IT", "Physics"]}],
        "singapore": [{"name": "National University of Singapore (NUS)", "ranking": "#8 QS 2024", "programs": ["Engineering", "Business", "Medicine"]}, {"name": "Nanyang Technological University (NTU)", "ranking": "#26 QS 2024", "programs": ["Engineering", "Business", "Computer Science"]}, {"name": "Singapore Management University", "ranking": "Top Business School", "programs": ["Business", "Law", "Information Systems"]}],
        "spain": [{"name": "University of Barcelona", "ranking": "#171 QS 2024", "programs": ["Medicine", "Law", "Arts"]}, {"name": "IE University", "ranking": "AACSB Accredited", "programs": ["Business", "Law", "Architecture"]}, {"name": "University of Madrid (Complutense)", "ranking": "Top 250 QS", "programs": ["Law", "Medicine", "Sciences"]}],
        "switzerland": [{"name": "ETH Zurich", "ranking": "#7 QS 2024", "programs": ["Engineering", "Computer Science", "Sciences"]}, {"name": "EPFL", "ranking": "#19 QS 2024", "programs": ["Engineering", "Life Sciences", "Computer Science"]}, {"name": "University of Zurich", "ranking": "Top 100 QS", "programs": ["Medicine", "Law", "Sciences"]}],
        "sweden": [{"name": "KTH Royal Institute of Technology", "ranking": "#89 QS 2024", "programs": ["Engineering", "Computer Science", "Architecture"]}, {"name": "Lund University", "ranking": "#92 QS 2024", "programs": ["Medicine", "Engineering", "Law"]}, {"name": "Stockholm University", "ranking": "Top 200 QS", "programs": ["Sciences", "Business", "Law"]}],
        "uk": [{"name": "University of Oxford", "ranking": "#3 QS 2024", "programs": ["Law", "Medicine", "Sciences"]}, {"name": "University of Cambridge", "ranking": "#2 QS 2024", "programs": ["Engineering", "Sciences", "Law"]}, {"name": "Imperial College London", "ranking": "#6 QS 2024", "programs": ["Engineering", "Medicine", "Sciences"]}, {"name": "University College London", "ranking": "#9 QS 2024", "programs": ["Medicine", "Engineering", "Arts"]}],
        "usa": [{"name": "Massachusetts Institute of Technology (MIT)", "ranking": "#1 QS 2024", "programs": ["Engineering", "Computer Science", "Sciences"]}, {"name": "Stanford University", "ranking": "#5 QS 2024", "programs": ["Business", "Engineering", "Computer Science"]}, {"name": "Harvard University", "ranking": "#4 QS 2024", "programs": ["Law", "Medicine", "Business"]}, {"name": "Carnegie Mellon University", "ranking": "#52 QS 2024", "programs": ["Computer Science", "Engineering", "Business"]}],
    }

    top_unis = top_unis_map.get(slug, [
        {"name": name + " National University", "ranking": "Top Ranked", "programs": programs[:3]},
    ])

    # Use extracted unis to augment if available
    if unis:
        for u in unis[:4]:
            existing_names = [t["name"].lower() for t in top_unis]
            if u.lower() not in existing_names:
                top_unis.append({"name": u, "ranking": "", "programs": programs[:2]})

    return {
        "slug": slug,
        "name": name,
        "flag": FLAGS.get(slug, "🌍"),
        "tagline": taglines.get(slug, "Study Abroad Destination"),
        "description": descriptions.get(slug, f"Discover world-class education opportunities in {name}."),
        "heroImage": UNSPLASH_HEROES.get(slug, ""),
        "images": images[:6],
        "stats": stats,
        "tuitionFees": tuition,
        "scholarships": scholarships if scholarships else [
            {"name": f"{name} Government Scholarship", "amount": "Varies", "eligibility": "Merit-based for international students", "link": ""},
            {"name": "University Excellence Awards", "amount": "Up to 50% tuition waiver", "eligibility": "Academic excellence; GPA 3.5+", "link": ""},
        ],
        "visaRequirements": visa,
        "workPermit": work_permit,
        "popularPrograms": programs,
        "languageRequirements": lang_req,
        "topUniversities": top_unis[:8],
        "faqs": faqs,
    }


# ── Main ───────────────────────────────────────────────────────────────────

def process_file(filename, slug):
    filepath = BLOGS_DIR / filename
    if not filepath.exists():
        print(f"  ⚠️  File not found: {filepath}")
        return None

    print(f"\n📄 Processing: {filename} → {slug}")
    ext = filepath.suffix.lower()

    try:
        if ext == ".pptx":
            text = extract_text_from_pptx(filepath)
            images = extract_images_from_pptx(filepath, slug)
        elif ext == ".pdf":
            text = extract_text_from_pdf(filepath)
            images = extract_images_from_pdf(filepath, slug)
        else:
            print(f"  ⚠️  Unknown extension: {ext}")
            return None

        print(f"  ✅ Text extracted: {len(text)} chars, {len(images)} images saved")
        data = build_country_json(slug, text, images)

        out_path = JSON_OUT_DIR / f"{slug}.json"
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)

        print(f"  💾 Saved: {out_path}")
        return data

    except Exception as e:
        print(f"  ❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return None


def main():
    print("=" * 60)
    print("MVR Blog Extractor — Country Data Generator")
    print("=" * 60)

    target = sys.argv[1] if len(sys.argv) > 1 else None
    results = {}

    for filename, slug in FILE_SLUG_MAP.items():
        if target and slug != target:
            continue
        data = process_file(filename, slug)
        if data:
            results[slug] = data

    print(f"\n{'=' * 60}")
    print(f"✅ Done! Processed {len(results)}/{len(FILE_SLUG_MAP)} countries")
    print(f"📁 JSON files saved to: {JSON_OUT_DIR}")
    print(f"🖼️  Images saved to: {IMG_OUT_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    main()
